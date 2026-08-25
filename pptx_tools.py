"""
📊 Taqdimot (PPTX) qurish — AI tomonidan JSON ko'rinishida tuzilgan slaydlar
ro'yxatini (sarlavha + bullet nuqtalar) chiroyli, izchil dizaynli PowerPoint
faylga aylantiradi. Har bir slaydda bitta rang sxemasi, bitta shrift oilasi
ishlatiladi — natija tasodifiy emas, balki "shablon" ko'rinishida chiqadi.
"""

from io import BytesIO

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ============================================================
# DIZAYN TOKENLARI — bitta joyda o'zgartirilsa, BUTUN taqdimot uslubi
# izchil o'zgaradi (professional taqdimotlarning asosiy belgisi).
# ============================================================
_BG = RGBColor(0x0F, 0x17, 0x2A)          # to'q ko'k-qora fon
_ACCENT = RGBColor(0x4C, 0x8B, 0xF5)       # och ko'k urg'u rangi (sarlavhalar, chiziqlar)
_ACCENT_SOFT = RGBColor(0x8A, 0xB4, 0xF8)  # yumshoqroq urg'u (bullet belgilari)
_TEXT_LIGHT = RGBColor(0xF2, 0xF4, 0xF8)   # asosiy matn (deyarli oq)
_TEXT_MUTED = RGBColor(0xB8, 0xC2, 0xD6)   # ikkinchi darajali matn (izoh, slayd raqami)
_FONT_TITLE = "Calibri"
_FONT_BODY = "Calibri"

_SLIDE_W = Inches(13.333)   # 16:9
_SLIDE_H = Inches(7.5)


def _set_solid_bg(slide, color: RGBColor):
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = color


def _add_text(shape, text: str, size: int, color: RGBColor, bold: bool = False,
               align=PP_ALIGN.LEFT, font=_FONT_BODY, italic: bool = False):
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font
    run.font.color.rgb = color
    return p


def _accent_bar(slide, x, y, w, h):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    bar.fill.solid()
    bar.fill.fore_color.rgb = _ACCENT
    bar.line.fill.background()
    bar.shadow.inherit = False
    return bar


def _slide_number(slide, current: int, total: int):
    box = slide.shapes.add_textbox(_SLIDE_W - Inches(1.6), _SLIDE_H - Inches(0.55), Inches(1.3), Inches(0.4))
    _add_text(box, f"{current} / {total}", 11, _TEXT_MUTED, align=PP_ALIGN.RIGHT)


def _title_slide(prs: Presentation, title: str, subtitle: str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # bo'sh layout
    _set_solid_bg(slide, _BG)
    _accent_bar(slide, Inches(0), Inches(3.55), Inches(1.4), Pt(6))

    title_box = slide.shapes.add_textbox(Inches(0.9), Inches(2.6), Inches(11.5), Inches(1.8))
    _add_text(title_box, title, 40, _TEXT_LIGHT, bold=True, font=_FONT_TITLE)

    sub_box = slide.shapes.add_textbox(Inches(0.95), Inches(4.0), Inches(11), Inches(0.9))
    _add_text(sub_box, subtitle, 18, _ACCENT_SOFT, font=_FONT_BODY)


def _content_slide(prs: Presentation, index: int, total: int, heading: str, bullets: list[str]):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_solid_bg(slide, _BG)
    _accent_bar(slide, Inches(0.9), Inches(1.15), Inches(0.9), Pt(5))

    head_box = slide.shapes.add_textbox(Inches(0.9), Inches(0.5), Inches(11.5), Inches(1.0))
    _add_text(head_box, heading, 28, _TEXT_LIGHT, bold=True, font=_FONT_TITLE)

    body = slide.shapes.add_textbox(Inches(1.1), Inches(1.7), Inches(11.0), Inches(5.2))
    tf = body.text_frame
    tf.word_wrap = True
    for i, point in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(14)
        run = p.add_run()
        run.text = f"●  {point}"
        run.font.size = Pt(19)
        run.font.name = _FONT_BODY
        run.font.color.rgb = _TEXT_LIGHT

    _slide_number(slide, index, total)


def _closing_slide(prs: Presentation, text: str, total: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_solid_bg(slide, _BG)
    _accent_bar(slide, Inches(0), Inches(3.55), Inches(1.4), Pt(6))
    box = slide.shapes.add_textbox(Inches(0.9), Inches(3.0), Inches(11.5), Inches(1.5))
    _add_text(box, text, 34, _TEXT_LIGHT, bold=True, font=_FONT_TITLE)
    _slide_number(slide, total, total)


def build_presentation(topic: str, slides: list[dict], author_note: str = "") -> BytesIO:
    """
    slides: [{"heading": str, "bullets": [str, ...]}, ...]
    Natija: sarlavha slaydi + har bir `slides` elementi uchun bitta kontent
    slayd + yakuniy "Eʼtiboringiz uchun rahmat!" slaydi.
    """
    prs = Presentation()
    prs.slide_width = _SLIDE_W
    prs.slide_height = _SLIDE_H

    total = len(slides) + 2  # titul + kontent slaydlar + yakuniy
    _title_slide(prs, topic, author_note or "Taqdimot")

    for i, s in enumerate(slides, start=2):
        _content_slide(prs, i, total, s.get("heading", f"{i - 1}-slayd"), s.get("bullets", []) or ["—"])

    _closing_slide(prs, "E'tiboringiz uchun rahmat!", total)

    buffer = BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer
